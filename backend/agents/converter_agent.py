"""
文档转化 Agent
1. 用 pymupdf/python-docx/python-pptx 提取原始文本（本地，无模型）
2. 将粗文本发送给 MiniMax M2.5 整理为结构化 Markdown
"""

import io
import re
import structlog
from services.model_router import model_router
from prompts.conversion import build_conversion_prompt

logger = structlog.get_logger()

SUPPORTED_FORMATS = {".doc", ".docx", ".pdf", ".ppt", ".pptx", ".xls", ".xlsx", ".csv", ".md", ".txt"}
# OLE2 二进制旧版 Office 格式 — 上传后先用 libreoffice headless 转成 OOXML 新格式
LEGACY_OFFICE_FORMATS = {".doc": ".docx", ".ppt": ".pptx", ".xls": ".xlsx"}
# 转换分段大小：~10K 字符约对应 5-7K tokens 输入，配合 max_tokens=8000 输出
# 留约 1.2x 膨胀空间和 reasoning model 的 <think> 预算，避免输出被截断。
# 历史值是 50000，配 8000 输出强制 4x 压缩，导致大量文档后半段丢失。
MAX_CHUNK_CHARS = 10000
# 截断检测阈值：单段输出字符数若 ≥ 该值，认为命中 max_tokens 上限被截断
TRUNCATION_WARN_CHARS = 11000


def _docx_zip_fallback(content: bytes) -> str:
    """绕开 python-docx 直接从 zip 抽 word/document.xml 的纯文本。

    用于以下场景:
    - python-docx 报 "There is no item named 'NULL' in the archive"
      (WPS / Pages / LibreOffice 转换出的 docx,Content_Types.xml 里 PartName 异常)
    - 文件 OOXML 结构基本完整但 manifest 不规范
    - 其他 python-docx 严格校验失败但 zip 本身可读

    实现:zip 里找 word/document.xml,解析所有 <w:t> 文本节点拼起来。
    丢失格式 / 表格结构,但保留文字内容(对 LLM 后续理解够用)。
    """
    import zipfile, re
    from xml.etree import ElementTree as ET

    with zipfile.ZipFile(io.BytesIO(content)) as z:
        names = z.namelist()
        # 优先 word/document.xml,兜底任何 word/document*.xml
        target = None
        if "word/document.xml" in names:
            target = "word/document.xml"
        else:
            for n in names:
                if n.startswith("word/document") and n.endswith(".xml"):
                    target = n; break
        if not target:
            raise RuntimeError(f"docx zip 里找不到 word/document.xml (zip entries: {names[:8]}...)")
        xml_bytes = z.read(target)

    # 解析所有 <w:t> 文本(忽略命名空间)
    # 用宽松的正则兜底,避免 XML 解析对损坏文件再次抛错
    try:
        root = ET.fromstring(xml_bytes)
        ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
        texts = [t.text for t in root.iter(f"{ns}t") if t.text]
        # <w:p> 段落分隔
        paras = []
        cur = []
        for elem in root.iter():
            tag = elem.tag.replace(ns, "")
            if tag == "p":
                if cur:
                    paras.append("".join(cur).strip())
                    cur = []
            elif tag == "t" and elem.text:
                cur.append(elem.text)
        if cur:
            paras.append("".join(cur).strip())
        return "\n".join(p for p in paras if p)
    except ET.ParseError:
        # XML 也炸了 — 用正则硬抽 <w:t>...</w:t>
        text = xml_bytes.decode("utf-8", errors="replace")
        chunks = re.findall(r'<w:t[^>]*>([^<]*)</w:t>', text)
        return "\n".join(c for c in chunks if c.strip())


def extract_text_from_docx(content: bytes) -> str:
    """优先用 python-docx 严格解析(保留段落 + 表格);
    任何异常(包括 lazy-load 时的 KeyError)→ fallback 到 zip 直读。

    典型 fallback 触发:
      - "There is no item named 'NULL' in the archive" (WPS / Pages 异常 manifest)
      - python-docx 找不到 [Content_Types].xml 里声明的 part
      - zip 损坏但 word/document.xml 仍可读
    """
    from docx import Document
    try:
        doc = Document(io.BytesIO(content))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        msg = str(e)[:120]
        logger.warning("docx_strict_parse_failed_fallback", err=msg, exc_type=type(e).__name__)
        try:
            text = _docx_zip_fallback(content)
            if text.strip():
                logger.info("docx_zip_fallback_ok", chars=len(text), prev_err=msg)
                return text
            logger.warning("docx_zip_fallback_empty", prev_err=msg)
        except Exception as fb_err:
            logger.warning("docx_zip_fallback_also_failed", err=str(fb_err)[:120])
        # 兜底也失败 → 抛原异常,让 task retry / 用户看错误
        raise


def extract_text_from_pdf(content: bytes) -> str:
    """优先抽 PDF 文本层(PyMuPDF);
    若总字数 < 平均每页 50 字 → 判定扫描件,日志提示走 vision OCR(由 convert_to_markdown 调度)。

    这里只返回文本层结果,扫描件检测的入口在 convert_to_markdown 里 — 因为 vision OCR 是异步,
    且要复用 model_router,放在 convert_to_markdown 同步路径里更顺。
    """
    import fitz
    doc = fitz.open(stream=content, filetype="pdf")
    parts = []
    for page_num, page in enumerate(doc, 1):
        text = page.get_text()
        if text.strip():
            parts.append(f"[第{page_num}页]\n{text}")
    return "\n\n".join(parts)


def _pdf_page_count(content: bytes) -> int:
    import fitz
    return fitz.open(stream=content, filetype="pdf").page_count


def _is_scanned_pdf(text: str, page_count: int) -> bool:
    """判定:平均每页文本 < 50 字 → 当扫描件处理。
    这个阈值能捕获完全无文本层的纯扫描件 + 文本层稀疏(只 OCR 了一部分)的混合件,
    可能误伤 PPT 转 PDF / 封面页 + 正文有文本但极简的文档,但召回率优先。
    """
    if page_count <= 0:
        return False
    avg = len(text.strip()) / page_count
    return avg < 50


XIAOMI_VISION_ENDPOINT = "https://api.xiaomimimo.com/v1/chat/completions"
XIAOMI_VISION_MODEL = "mimo-v2-omni"      # 或 mimo-v2.5,均支持 vision


async def _resolve_xiaomi_key() -> str:
    """读 xiaomi_api_key:DB(config_service api_keys 表)优先,fallback .env settings。"""
    try:
        from services.config_service import config_service
        db_entry = await config_service.get("api_keys", "xiaomi_api_key")
        if db_entry and db_entry.get("value"):
            return db_entry["value"]
    except Exception:
        pass
    from config import settings
    return getattr(settings, "xiaomi_api_key", "") or ""


async def _call_xiaomi_vision_one_page(
    *, page_idx: int, n_pages: int, b64: str, api_key: str, client,
) -> tuple[int, str]:
    """按官方文档形态直接调 https://api.xiaomimimo.com/v1/chat/completions。

    与 model_router.chat 的差异:
    - endpoint 不同(api.xiaomimimo.com vs token-plan-cn 代理)
    - header 用 api-key 而不是 Authorization: Bearer
    - body 用 max_completion_tokens 而不是 max_tokens

    单页失败返回错误占位文本,不抛异常 — 让整篇 OCR 不被一页拖垮。
    """
    import httpx
    payload = {
        "model": XIAOMI_VISION_MODEL,
        "messages": [{
            "role": "user",
            "content": [
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{b64}"},
                },
                {
                    "type": "text",
                    "text": (
                        f"这是一份 PDF 扫描件的第 {page_idx + 1}/{n_pages} 页。请忠实转写本页所有文字内容为 markdown:\n"
                        "- 保留段落、标题层级(用 # ## ###)、表格(用 markdown table 语法)、列表(- / 1.)\n"
                        "- 印章 / 签名 / 手写批注用 [印章: ...]、[签名: ...]、[手写: ...] 标注位置\n"
                        "- 公式 / 图表里的数据要转写,纯装饰性图忽略\n"
                        "- **不要总结,不要补全或修正你认为应该有的内容**,只忠实转写图上看到的文字\n"
                        "- 整页空白或仅装饰图无文字 → 输出 \"(本页无文字)\"\n"
                        "- 不要包裹 ```markdown 代码块,直接输出 markdown 正文"
                    ),
                },
            ],
        }],
        "max_completion_tokens": 4000,
        "temperature": 0.1,
    }
    try:
        r = await client.post(
            XIAOMI_VISION_ENDPOINT,
            headers={"api-key": api_key, "Content-Type": "application/json"},
            json=payload,
            timeout=180.0,
        )
        if r.status_code != 200:
            logger.warning("xiaomi_vision_http_error", page=page_idx + 1,
                           status=r.status_code, body=r.text[:200])
            return page_idx, f"(第 {page_idx + 1} 页识别失败: HTTP {r.status_code})"
        data = r.json()
        text = (data.get("choices") or [{}])[0].get("message", {}).get("content") or ""
        return page_idx, text.strip()
    except Exception as e:
        logger.warning("xiaomi_vision_exception", page=page_idx + 1, err=str(e)[:120])
        return page_idx, f"(第 {page_idx + 1} 页识别失败: {str(e)[:80]})"


async def _ocr_pdf_via_vision_llm(content: bytes) -> str:
    """扫描件 PDF → 多模态 LLM 逐页转写 markdown。

    流程:
      1. PyMuPDF 把每页 render 为 PNG (dpi=180,平衡清晰度 + base64 体积)
      2. base64 编码塞 OpenAI 多模态格式 messages[].content[].image_url
      3. 并发 5 路调 https://api.xiaomimimo.com/v1/chat/completions (mimo-v2-omni)
      4. 按页码拼回 markdown,失败页标占位不阻断整篇

    走官方 endpoint + api-key header(不是 model_router 的 token-plan-cn 代理),
    因为 vision 调用按官方文档形态最稳。

    成本权衡:每页 1 次调用 ≈ 几千 token (图像编码后)。30 页 PDF 约 ¥0.5-2。
    """
    import asyncio, base64
    import httpx
    import fitz

    api_key = await _resolve_xiaomi_key()
    if not api_key:
        logger.warning("xiaomi_vision_no_key")
        return "(扫描件 PDF — 但未配置 xiaomi_api_key,无法走 OCR)"

    doc = fitz.open(stream=content, filetype="pdf")
    n_pages = doc.page_count
    if n_pages == 0:
        doc.close()
        return ""

    # render 所有页为 PNG → base64
    pages_b64: list[str] = []
    for page in doc:
        pix = page.get_pixmap(dpi=180)
        png = pix.tobytes("png")
        pages_b64.append(base64.b64encode(png).decode("ascii"))
    doc.close()

    sem = asyncio.Semaphore(5)

    async with httpx.AsyncClient(timeout=180.0) as client:
        async def _ocr_one(idx: int, b64: str) -> tuple[int, str]:
            async with sem:
                return await _call_xiaomi_vision_one_page(
                    page_idx=idx, n_pages=n_pages, b64=b64,
                    api_key=api_key, client=client,
                )
        results = await asyncio.gather(*[_ocr_one(i, b) for i, b in enumerate(pages_b64)])

    results.sort(key=lambda x: x[0])
    blocks: list[str] = []
    for i, md in results:
        blocks.append(f"## 第 {i + 1} 页\n\n{md}".rstrip())
    return "\n\n---\n\n".join(blocks)


def extract_text_from_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        if slide_texts:
            parts.append(f"[幻灯片{i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def extract_text_from_xlsx(content: bytes) -> str:
    import openpyxl

    def _extract_openpyxl(read_only: bool) -> str:
        wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True, read_only=read_only)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(v) if v is not None else "" for v in row)
                if row_text.strip(" |"):
                    parts.append(row_text)
        return "\n".join(parts)

    def _extract_zip_fallback() -> str:
        """最后兜底：从 xlsx zip 中直接解析 sharedStrings.xml 提取所有文本。"""
        import zipfile
        import xml.etree.ElementTree as ET
        parts = []
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as z:
                if "xl/sharedStrings.xml" in z.namelist():
                    with z.open("xl/sharedStrings.xml") as f:
                        tree = ET.parse(f)
                        ns = {"x": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
                        for si in tree.findall(".//x:si", ns):
                            texts = [t.text for t in si.findall(".//x:t", ns) if t.text]
                            val = "".join(texts).strip()
                            if val:
                                parts.append(val)
        except Exception:
            pass
        return "\n".join(parts) if parts else "[无法解析此xlsx文件]"

    # 1. 普通模式
    try:
        return _extract_openpyxl(read_only=False)
    except Exception:
        pass
    # 2. 流式只读模式（跳过 DataValidation / 样式解析）
    try:
        return _extract_openpyxl(read_only=True)
    except Exception:
        pass
    # 3. 直接从 zip 解析 sharedStrings（兜底，无格式但不丢文本）
    return _extract_zip_fallback()


def _soffice_convert(content: bytes, src_ext: str, target_ext: str) -> bytes:
    """用 libreoffice headless 把 .doc/.ppt/.xls 转成 .docx/.pptx/.xlsx。

    target_ext 不含点 (传 'docx' / 'pptx' / 'xlsx')。
    soffice 必须在 PATH 里(Dockerfile 装了 libreoffice 包)。
    """
    import subprocess, tempfile, os
    target_clean = target_ext.lstrip(".")
    with tempfile.TemporaryDirectory() as tmpdir:
        in_path = os.path.join(tmpdir, f"input{src_ext}")
        with open(in_path, "wb") as f:
            f.write(content)
        # soffice --headless --convert-to docx input.doc --outdir /tmp/x
        try:
            proc = subprocess.run(
                ["soffice", "--headless", "--convert-to", target_clean,
                 "--outdir", tmpdir, in_path],
                capture_output=True, timeout=120,
            )
        except FileNotFoundError:
            raise RuntimeError("libreoffice 未安装(需 apt-get install libreoffice)")
        except subprocess.TimeoutExpired:
            raise RuntimeError(f"libreoffice 转换超时 (>120s) — 文件可能损坏或过大")
        if proc.returncode != 0:
            raise RuntimeError(
                f"libreoffice 转换失败 (rc={proc.returncode}): {proc.stderr.decode('utf-8', errors='replace')[:200]}"
            )
        # 找输出文件 (input.docx 等)
        out_name = f"input.{target_clean}"
        out_path = os.path.join(tmpdir, out_name)
        if not os.path.exists(out_path):
            # 兜底:扫目录里第一个匹配后缀的文件
            for fn in os.listdir(tmpdir):
                if fn.endswith(f".{target_clean}"):
                    out_path = os.path.join(tmpdir, fn)
                    break
            else:
                raise RuntimeError(f"libreoffice 没生成 {target_clean} 文件")
        with open(out_path, "rb") as f:
            return f.read()


def _normalize_legacy_office(filename: str, content: bytes) -> tuple[str, bytes]:
    """如果是 .doc/.ppt/.xls,先用 soffice 转成 .docx/.pptx/.xlsx,返回新文件名 + 内容。
    其他格式原样返回。
    """
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext in LEGACY_OFFICE_FORMATS:
        target_ext = LEGACY_OFFICE_FORMATS[ext]
        logger.info("legacy_office_convert", src=ext, target=target_ext, filename=filename)
        new_content = _soffice_convert(content, src_ext=ext, target_ext=target_ext)
        # 替换扩展名
        base = filename.rsplit(".", 1)[0] if "." in filename else filename
        return f"{base}{target_ext}", new_content
    return filename, content


def extract_raw_text(filename: str, content: bytes) -> str:
    # 老 Office 格式先转换成新格式,再走原路径
    filename, content = _normalize_legacy_office(filename, content)
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == ".docx":
        return extract_text_from_docx(content)
    elif ext == ".pdf":
        return extract_text_from_pdf(content)
    elif ext == ".pptx":
        return extract_text_from_pptx(content)
    elif ext in (".xlsx",):
        return extract_text_from_xlsx(content)
    elif ext in (".csv", ".md", ".txt"):
        return content.decode("utf-8", errors="replace")
    else:
        raise ValueError(f"不支持的格式: {ext}")


def _xlsx_raw_to_markdown(raw_text: str, doc_title: str) -> str:
    """xlsx/csv 抽取出来的原始文本本身已结构化（[Sheet:xx] + " | " 分隔的行），
    转 markdown 不需要 LLM 重新理解，纯本地处理：
      - [Sheet: xx] -> ## xx
      - 每个 sheet 第一行作为表头，紧跟 | --- | 分隔行
      - 数据行用 markdown table 格式 | a | b | c |
      - cell 内换行替换为 <br>
    """
    out: list[str] = [f"# {doc_title}", ""]
    lines = [ln for ln in raw_text.split("\n") if ln.strip()]
    in_table = False
    is_first_row = True
    col_count = 0

    def _normalize_row(line: str) -> tuple[str, int]:
        cells = [c.strip().replace("\n", "<br>") for c in line.split(" | ")]
        return "| " + " | ".join(cells) + " |", len(cells)

    for ln in lines:
        if ln.startswith("[Sheet: ") and ln.endswith("]"):
            sheet_name = ln[len("[Sheet: "):-1].strip()
            out.append("")
            out.append(f"## Sheet：{sheet_name}")
            out.append("")
            in_table = False
            is_first_row = True
            continue

        # 普通文本行（无 ` | `）
        if " | " not in ln:
            # 收尾上一段表格
            if in_table:
                out.append("")
                in_table = False
                is_first_row = True
            out.append(ln)
            continue

        # 表格行
        row_md, cnt = _normalize_row(ln)
        if is_first_row or cnt != col_count:
            # 起新表 / 列数变化时另起一表
            if in_table and not is_first_row:
                out.append("")
            out.append(row_md)
            out.append("| " + " | ".join(["---"] * cnt) + " |")
            col_count = cnt
            in_table = True
            is_first_row = False
        else:
            out.append(row_md)

    return "\n".join(out).strip() + "\n"


async def convert_to_markdown(filename: str, content: bytes) -> tuple[str, str | None]:
    """Returns (markdown_text, model_name_or_None)."""
    # 老 Office 格式 (.doc/.ppt/.xls) 先用 soffice 转成新格式,后续 ext 判断按新格式走
    filename, content = _normalize_legacy_office(filename, content)
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext in (".md", ".txt"):
        return content.decode("utf-8", errors="replace"), None

    # PDF:先抽文本层,若是扫描件(平均每页<50字)→ 走多模态 vision OCR
    if ext == ".pdf":
        text_layer = extract_raw_text(filename, content)
        try:
            n_pages = _pdf_page_count(content)
        except Exception:
            n_pages = 0
        if _is_scanned_pdf(text_layer, n_pages):
            logger.info(
                "pdf_scanned_detected_using_vision_ocr",
                filename=filename, pages=n_pages, text_layer_chars=len(text_layer.strip()),
            )
            ocr_md = await _ocr_pdf_via_vision_llm(content)
            return ocr_md, "mimo-v2-omni"
        # 文本层够 → 走原 LLM 整理路径
        raw_text = text_layer
        logger.info("text_extracted", filename=filename, chars=len(raw_text), source="text_layer")
    else:
        raw_text = extract_raw_text(filename, content)
        logger.info("text_extracted", filename=filename, chars=len(raw_text))

    # xlsx / csv 走本地 fast path：原始抽取已是结构化表格，LLM 能加的价值有限，
    # 直接构造 markdown 既零 token 消耗，也不会被 max_tokens 截断。
    if ext in (".xlsx", ".csv"):
        if not raw_text.strip():
            return "", None
        doc_title = filename.rsplit(".", 1)[0]
        md = _xlsx_raw_to_markdown(raw_text, doc_title)
        logger.info("converted_local_fast_path", filename=filename, ext=ext, chars=len(md))
        return md, None

    # 分段发送（超过 MAX_CHUNK_CHARS 时）
    if len(raw_text) <= MAX_CHUNK_CHARS:
        segments = [raw_text]
    else:
        segments = [raw_text[i:i + MAX_CHUNK_CHARS] for i in range(0, len(raw_text), MAX_CHUNK_CHARS)]

    # 分段并发调 LLM —— Semaphore=3 限流，避免 provider 限速
    import asyncio
    sem = asyncio.Semaphore(3)

    async def _convert_one(idx: int, segment: str) -> tuple[int, str, str | None]:
        async with sem:
            prompt = await build_conversion_prompt(segment)
            result, model = await model_router.chat_with_routing(
                "conversion",
                [{"role": "user", "content": prompt}],
                max_tokens=8000,
                timeout=180.0,
            )
            result = re.sub(r"<think>[\s\S]*?</think>", "", result, flags=re.IGNORECASE).strip()
            # 截断检测：输出字符数贴近 max_tokens 对应字符上限就 warn
            if len(result) >= TRUNCATION_WARN_CHARS:
                logger.warning(
                    "segment_likely_truncated",
                    filename=filename, segment=idx + 1, total=len(segments),
                    output_chars=len(result), input_chars=len(segment), model=model,
                )
            logger.info("segment_converted", filename=filename, segment=idx + 1, total=len(segments), model=model)
            return idx, result, model

    outputs = await asyncio.gather(*[_convert_one(i, s) for i, s in enumerate(segments)])
    outputs.sort(key=lambda x: x[0])
    markdown_parts = [o[1] for o in outputs]
    used_model = next((o[2] for o in outputs if o[2]), None)

    return "\n\n".join(markdown_parts), used_model
