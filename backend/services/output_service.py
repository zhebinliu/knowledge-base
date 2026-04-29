"""Core generation logic for CuratedBundle outputs (kickoff_pptx / survey / insight).

对话式访谈流程：CuratedBundle.extra.conversation_id 指向 OutputConversation。
生成器读取完整对话记录 + 访谈过程中检索到的 refs，作为主要素材；
project_id 可能为 None（行业作用域）。
"""
import io
import structlog
from datetime import date
from config import settings
from models import async_session_maker
from models.curated_bundle import CuratedBundle
from models.project import Project
from models.output_conversation import OutputConversation
from models.agent_config import AgentConfig
from sqlalchemy import select

logger = structlog.get_logger()


async def _get_project(project_id: str) -> Project | None:
    if not project_id:
        return None
    async with async_session_maker() as s:
        return await s.get(Project, project_id)


async def _get_conversation(conv_id: str) -> OutputConversation | None:
    if not conv_id:
        return None
    async with async_session_maker() as s:
        return await s.get(OutputConversation, conv_id)


async def _get_output_agent_config(key: str) -> dict:
    async with async_session_maker() as s:
        row = (await s.execute(
            select(AgentConfig).where(
                AgentConfig.config_type == "output_agent",
                AgentConfig.config_key == key,
            )
        )).scalar_one_or_none()
    if row and isinstance(row.config_value, dict):
        return {
            "prompt": row.config_value.get("prompt", ""),
            "skill_ids": row.config_value.get("skill_ids", []),
            "model": row.config_value.get("model"),
        }
    return {"prompt": "", "skill_ids": [], "model": None}


async def _get_skill_snippets(skill_ids: list[str]) -> str:
    if not skill_ids:
        return ""
    from models.skill import Skill
    async with async_session_maker() as s:
        rows = (await s.execute(select(Skill).where(Skill.id.in_(skill_ids)))).scalars().all()
    return "\n\n".join(f"### 技能：{r.name}\n{r.prompt_snippet}" for r in rows)


def _format_transcript(conv: OutputConversation | None) -> str:
    if not conv or not conv.messages:
        return "（没有可用的访谈记录）"
    lines = []
    for m in conv.messages:
        role = m.get("role")
        if role == "user":
            content = (m.get("content") or "").strip()
            if not content or content.startswith("（请开场"):
                continue
            lines.append(f"**用户**：{content}")
        elif role == "assistant":
            content = (m.get("content") or "").strip()
            if content:
                lines.append(f"**顾问**：{content}")
    return "\n\n".join(lines) if lines else "（没有可用的访谈记录）"


def _format_refs(conv: OutputConversation | None) -> str:
    if not conv or not conv.refs:
        return ""
    # 去重（按 chunk_id）
    seen = set()
    lines = []
    for r in conv.refs:
        cid = r.get("chunk_id")
        if cid in seen:
            continue
        seen.add(cid)
        header = f"[{r.get('filename') or '未知文档'}" + (f" · {r['source_section']}" if r.get("source_section") else "") + "]"
        lines.append(f"{header}\n{(r.get('content') or '')[:400]}")
    return "\n\n".join(lines[:20])


async def _llm_call(prompt: str, system: str = "", model: str | None = None, max_tokens: int | None = 8000, timeout: float = 180.0) -> str:
    from services.model_router import model_router
    messages = []
    if system:
        messages.append({"role": "system", "content": system})
    messages.append({"role": "user", "content": prompt})
    if model:
        content, _ = await model_router.chat(model, messages, max_tokens=max_tokens, timeout=timeout)
    else:
        content, _ = await model_router.chat_with_routing("doc_generation", messages, max_tokens=max_tokens, timeout=timeout)
    return content


async def _mark_bundle(bundle_id: str, status: str, **kwargs):
    async with async_session_maker() as s:
        b = await s.get(CuratedBundle, bundle_id)
        if b:
            b.status = status
            for k, v in kwargs.items():
                setattr(b, k, v)
            await s.commit()


async def _mark_conversation(bundle_id: str, status: str):
    async with async_session_maker() as s:
        bundle = await s.get(CuratedBundle, bundle_id)
        if not bundle or not bundle.extra:
            return
        conv_id = bundle.extra.get("conversation_id")
        if not conv_id:
            return
        conv = await s.get(OutputConversation, conv_id)
        if conv:
            conv.status = status
            await s.commit()


INDUSTRY_PRIMING_SYSTEM = """你是一位资深的企业级 SaaS 实施研究员（前 Gartner / IDC 分析师背景）。
基于客户、行业、模块信息，输出一份结构化【行业/客户研究简报】，给到下游的咨询顾问做参考。
要求：
1. 不要虚构具体数字，但可以使用业内公开常识和典型范围（标"行业典型"或"业界基准"）
2. 不要写成营销话术，要写成内部研究备忘录的语气：客观、克制、能落地
3. 引用具体竞品 / 标杆客户名字时必须确认是真实存在的；不确定的写"同行业头部企业"
4. 用 Markdown 输出，包含下面 6 个块：
   ## 行业宏观
   - 当前规模 / 增速 / 政策动向 / 数字化渗透率（用业内公开口径）
   ## 客户画像
   - 该客户在行业里的位置（头部/中部/长尾）、典型业务模式、CRM 成熟度推断
   ## 同行业 CRM 实施常见模式
   - 头部客户怎么落 LTC / SFA / CPQ / Service？典型周期、典型团队规模、典型预算区间
   ## 该行业 CRM 实施的 5 大常见陷阱
   - 每条：陷阱名 / 触发条件 / 后果 / 规避建议
   ## 监管与合规要点
   - 行业相关的数据合规（个保法 / 行业专项 / 跨境数据）
   ## 关键成功要素（CSF）
   - 5–7 条，每条一句话结论
"""


async def _industry_priming(proj, industry: str | None, kind: str, model: str | None) -> str:
    """让 LLM 当行业研究员，先生成一份结构化的行业/客户研究 brief，作为后续生成的素材。
    这相当于用模型自己的训练知识做一轮"准检索"。"""
    if not industry and not (proj and proj.customer):
        return ""
    customer = proj.customer if proj else ""
    modules = ", ".join(proj.modules or []) if proj else ""
    prompt = f"""客户：{customer or "未提供"}
行业：{industry or "未提供"}
拟实施模块：{modules or "未提供"}
当前任务类型：{kind}

请按系统提示要求，输出该客户/行业的研究简报（800–1500 字）。"""
    try:
        return await _llm_call(prompt, system=INDUSTRY_PRIMING_SYSTEM, model=model, max_tokens=3000)
    except Exception as e:
        logger.warning("industry_priming_failed", err=str(e)[:120])
        return ""


async def _web_research(proj, industry: str | None, kind: str) -> tuple[list[dict], list[dict]]:
    """真实联网检索：返回 (条目列表, 调用日志)。没配 key 时返回空。"""
    from services.web_search_service import web_search, has_web_search_provider
    if not await has_web_search_provider():
        return [], []
    queries: list[str] = []
    if industry:
        queries.append(f"{industry} CRM 实施 案例 2024")
        queries.append(f"{industry} 数字化转型 痛点")
        if kind in ("kickoff_pptx", "kickoff_html"):
            queries.append(f"{industry} 龙头企业 销售管理 流程")
        elif kind == "insight":
            queries.append(f"{industry} CRM 项目 失败原因")
            queries.append(f"{industry} 客户管理 行业基准")
        elif kind == "survey":
            queries.append(f"{industry} 销售 调研 关键指标")
    if proj and proj.customer:
        queries.append(f"{proj.customer} 业务模式 数字化")
    queries = list(dict.fromkeys([q.strip() for q in queries if q.strip()]))[:5]

    items: list[dict] = []
    log: list[dict] = []
    for q in queries:
        hits = await web_search(q, top_k=5)
        log.append({"query": q, "hits": len(hits)})
        items.extend(hits)
    return items, log


def _format_web_items(items: list[dict]) -> str:
    if not items:
        return ""
    seen = set()
    blocks = []
    for it in items:
        u = it.get("url", "")
        if u in seen or not u:
            continue
        seen.add(u)
        title = it.get("title") or "—"
        snippet = (it.get("snippet") or "").strip()
        blocks.append(f"[{title}]({u})\n{snippet}")
        if len(blocks) >= 12:
            break
    return "\n\n".join(blocks)


async def _get_brief_block(project_id: str, kind: str) -> str:
    """读取已确认 Brief，渲染为 markdown 块；没有就返回空串。"""
    if not project_id:
        return ""
    from models.project_brief import ProjectBrief
    from services.brief_service import get_schema, render_brief_for_prompt
    schema = get_schema(kind)
    if not schema:
        return ""
    async with async_session_maker() as s:
        row = (await s.execute(
            select(ProjectBrief).where(
                ProjectBrief.project_id == project_id,
                ProjectBrief.output_kind == kind,
            )
        )).scalar_one_or_none()
    if not row or not row.fields:
        return ""
    return render_brief_for_prompt(row.fields, schema)


async def _gather_inputs(bundle_id: str, project_id: str, kind: str) -> dict:
    """统一拉取对话 / 项目 / 智能体配置，并在生成阶段做额外 KB 检索丰富上下文。"""
    async with async_session_maker() as s:
        bundle = await s.get(CuratedBundle, bundle_id)
    conv_id = (bundle.extra or {}).get("conversation_id") if bundle and bundle.extra else None
    industry_override = (bundle.extra or {}).get("industry") if bundle and bundle.extra else None
    conv = await _get_conversation(conv_id) if conv_id else None
    proj = await _get_project(project_id) if project_id else None
    agent_cfg = await _get_output_agent_config(kind)
    skill_text = await _get_skill_snippets(agent_cfg["skill_ids"])
    transcript = _format_transcript(conv)
    industry = (proj.industry if proj else None) or industry_override

    # 生成阶段额外检索：基于项目元数据 + kind 跑几轮 search_kb，把更多素材注入
    extra_refs, kb_log = await _generation_kb_search(proj, industry, kind, conv)
    all_refs_text = _merge_refs_text(_format_refs(conv), extra_refs)

    # 行业 / 客户研究 priming（使用模型自身知识做"准检索"）
    industry_brief = await _industry_priming(proj, industry, kind, agent_cfg.get("model"))
    # 项目 Brief（用户已确认/编辑过的字段，作为最权威的素材）
    project_brief_text = await _get_brief_block(project_id, kind)
    # 真实 web 搜索（仅当配置了 key）
    web_items, web_log = await _web_research(proj, industry, kind)
    web_text = _format_web_items(web_items)

    # 把所有检索日志写回 bundle.extra，前端可看
    async with async_session_maker() as s:
        b = await s.get(CuratedBundle, bundle_id)
        if b:
            b.extra = {
                **(b.extra or {}),
                "generation_kb_calls": kb_log,
                "web_search_calls": web_log,
                "has_industry_brief": bool(industry_brief),
            }
            await s.commit()
    return {
        "project": proj,
        "industry": industry,
        "conv": conv,
        "agent_prompt": agent_cfg["prompt"],
        "agent_model": agent_cfg["model"],
        "skill_text": skill_text,
        "transcript": transcript,
        "refs_text": all_refs_text,
        "industry_brief": industry_brief,
        "project_brief": project_brief_text,
        "web_text": web_text,
        "kb_calls": kb_log,
        "web_calls": web_log,
    }


async def _generation_kb_search(proj, industry: str | None, kind: str, conv) -> tuple[list[dict], list[dict]]:
    """生成阶段对症检索：按 kind 选 query 模板。返回 (额外 refs, [{query, hits}] 日志)。"""
    from agents.output_chat import _run_search_kb
    from models.document import Document

    queries: list[str] = []
    if proj:
        if proj.customer:
            queries.append(f"{proj.customer} {proj.industry or ''} 项目背景 关键挑战")
        for m in (proj.modules or [])[:3]:
            queries.append(f"{proj.industry or ''} {m} 实施要点 最佳实践")
    if industry:
        if kind in ("kickoff_pptx", "kickoff_html"):
            queries.extend([
                f"{industry} CRM 启动会 议程 范围",
                f"{industry} 实施 风险 应对",
                f"{industry} 团队 治理 RACI",
            ])
        elif kind == "insight":
            queries.extend([
                f"{industry} CRM 项目 风险 量化 指标",
                f"{industry} 干系人 决策点",
                f"{industry} 实施成功要素 经验",
            ])
        elif kind == "survey":
            queries.extend([
                f"{industry} 业务流程 调研 关键问题",
                f"{industry} 数据集成 接口 调研",
                f"{industry} 角色权限 合规要求",
            ])

    # 去重
    seen_q: set[str] = set()
    unique_queries = []
    for q in queries:
        q = q.strip()
        if q and q not in seen_q:
            seen_q.add(q)
            unique_queries.append(q)
    unique_queries = unique_queries[:6]

    project_document_ids: list[str] | None = None
    if proj:
        async with async_session_maker() as s:
            rows = (await s.execute(select(Document.id).where(Document.project_id == proj.id))).all()
        project_document_ids = [r[0] for r in rows] or None

    all_refs: list[dict] = []
    log: list[dict] = []
    for q in unique_queries:
        try:
            _text, refs = await _run_search_kb(q, project_document_ids, industry)
            log.append({"query": q, "hits": len(refs)})
            all_refs.extend(refs)
        except Exception as e:
            log.append({"query": q, "hits": 0, "error": str(e)[:80]})
    return all_refs, log


def _merge_refs_text(existing_text: str, extra_refs: list[dict]) -> str:
    """把对话期 refs 与生成期 refs 合并，按 chunk_id 去重，最多 30 条。"""
    blocks: list[str] = []
    if existing_text:
        blocks.append(existing_text)
    seen: set[str] = set()
    new_blocks: list[str] = []
    for r in extra_refs:
        cid = r.get("chunk_id")
        if not cid or cid in seen:
            continue
        seen.add(cid)
        header = f"[{r.get('filename') or '未知文档'}" + (f" · {r['source_section']}" if r.get("source_section") else "") + "]"
        new_blocks.append(f"{header}\n{(r.get('content') or '')[:400]}")
        if len(new_blocks) >= 15:
            break
    if new_blocks:
        blocks.append("【生成阶段补充检索】\n" + "\n\n".join(new_blocks))
    return "\n\n".join(blocks)


# ── Survey ────────────────────────────────────────────────────────────────────

SURVEY_SYSTEM = """你是 MBB 风格的资深 CRM 实施顾问（MECE 思维 / 金字塔原理）。
你正在为客户设计一份【实施前调研问卷】，交付给客户的项目经理 + 业务负责人填答。

【输出风格 — 严格遵守】
- 咨询公司内部文档体例：编号清晰、可勾选、可批注
- 不写"赋能/抓手/闭环/链路/生态"等黑话
- 每个问题要具体到可作答的颗粒度（"贵司销售从线索到签单的平均周期是多少天？" 而不是"贵司的销售周期如何？"）
- 区分【事实型】（一定有标准答案）/【判断型】（需主观评估）/【数据型】（需要从系统导出数据）
- 每题后面用斜体注明：*为什么问 / 答案如何使用*

【题量规模 — 至少要这么多，否则视为不合格】
- 总计 ≥ 60 题
- 7 个大类，每类 8–12 题：
  1. 战略与目标
  2. 组织与角色
  3. 业务流程（线索→机会→合同→交付→回款）
  4. 数据治理与主数据
  5. 系统集成与接口
  6. 合规、安全、权限
  7. 资源、预算与进度

【必须在问卷顶部加一段说明】
> 本问卷用于 CRM 实施启动前的现状摸底。请由对应模块责任人填写。带 ★ 的题目为重点题，请务必填答。

输出 Markdown 格式。"""


async def generate_survey(bundle_id: str, project_id: str):
    try:
        await _mark_bundle(bundle_id, "generating")
        ctx = await _gather_inputs(bundle_id, project_id, "survey")
        proj = ctx["project"]
        scope_line = f"项目：{proj.name}，客户：{proj.customer or '—'}" if proj else f"行业：{ctx['industry'] or '—'}"

        prompt = f"""{scope_line}
行业：{ctx['industry'] or '未填写'}

【访谈记录（已知信息）】
{ctx['transcript']}

{f"【行业/客户研究 brief（模型自身知识）】{chr(10)}{ctx['industry_brief']}" if ctx.get('industry_brief') else ""}

{f"【知识库佐证】{chr(10)}{ctx['refs_text']}" if ctx['refs_text'] else ""}

{f"【联网检索结果】{chr(10)}{ctx['web_text']}" if ctx.get('web_text') else ""}

{f"【方法论/风格要求】{chr(10)}{ctx['agent_prompt']}" if ctx['agent_prompt'] else ""}

{f"【启用技能】{chr(10)}{ctx['skill_text']}" if ctx['skill_text'] else ""}

请按系统提示要求，输出 ≥60 题的实施前调研问卷（7 个大类 × 8–12 题）。访谈已经有的信息**不要再问**，要根据已有信息**精准追问没覆盖的维度**。"""

        md = await _llm_call(prompt, system=SURVEY_SYSTEM, model=ctx["agent_model"], max_tokens=8000)

        docx_key: str | None = None
        try:
            title = f"调研问卷 · {proj.name if proj else (ctx['industry'] or '—')}"
            docx_bytes = _build_docx(title, md)
            docx_key = f"outputs/{bundle_id}/survey.docx"
            _minio_put(docx_key, docx_bytes, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        except Exception as e:
            logger.warning("survey_docx_failed", error=str(e)[:100])

        await _mark_bundle(bundle_id, "done", content_md=md, file_key=docx_key)
        await _mark_conversation(bundle_id, "done")
        logger.info("survey_generated", bundle_id=bundle_id, project_id=project_id)
    except Exception as e:
        logger.error("survey_failed", bundle_id=bundle_id, error=str(e)[:200])
        await _mark_bundle(bundle_id, "failed", error=str(e)[:500])
        await _mark_conversation(bundle_id, "failed")


# ── Insight ───────────────────────────────────────────────────────────────────

INSIGHT_SYSTEM = """你是 MBB 风格的资深咨询顾问（McKinsey / BCG / Bain），现在为客户高管层撰写【项目洞察报告】。

【风格 — 严格执行】
1. **金字塔原理**：每节先给"Bottom line"结论（1 句话，加粗），再展开论据
2. **Claim → Evidence → So what**：每个论点都要有数据支撑（行业 brief、知识库或访谈），并写出"对项目意味着什么"
3. **表格优先于 bullet**：能用表格表达的不要用 bullet（风险矩阵、决策表、干系人画像、量化指标）
4. **不写黑话**：禁止"赋能 / 抓手 / 闭环 / 链路 / 生态 / 数字化转型 / 一站式"
5. **数字化**：每节至少一个量化指标（区间也可以，标"业界基准 / 行业典型"）
6. **专业克制**：不要感叹号，不要营销话术，"我们认为"、"建议"用得克制
7. **可信度标注**：信息源用 [访谈] / [知识库] / [行业 brief] / [Web] 四种标签清晰区分；模型推断的用 [推断] 并说明依据
8. **不编造**：访谈和素材都没覆盖的，写"信息缺失，建议在 Phase 1 第一周补访"

【格式】
- 用 Markdown
- 表格用 GFM 语法
- 不要用 emoji
- 章节之间用 `---` 分隔
- 风险用「高/中/低」标签，不用 🔴🟡🟢

【关键约束 — 段落级输出】
- 你只负责输出**当前一节**的正文内容，不要再次输出整篇报告的标题、封面、元信息表（"客户名称 / 项目代号 / 报告日期 / 编写人 / 版本"等）
- **禁止**输出 `# 项目洞察报告` / `## 执行摘要` / `## {章节名}` 这类 H1/H2 标题，章节标题已由系统注入
- 直接从 Bottom line 加粗结论开始，需要更细粒度可用 H3 (`###`) 及以下
"""

INSIGHT_SECTIONS = [
    ("执行摘要", "用 3–5 条 bullet 概括项目当前态势、最关键 1 个机会与最关键 1 个风险，给高管 30 秒就能读完。"),
    ("项目概览", "基于访谈与知识库素材，梳理项目目标、范围、行业背景、当前阶段。要明确量化（用户数 / 模块数 / 预算 / 时间窗）。"),
    ("干系人画像", "甲方 / 乙方 / 联合团队的关键角色、决策权重、对项目的态度（积极 / 观望 / 阻力）。可以用表格列出。"),
    ("关键决策点", "梳理已做出的关键决策（含背景与影响）和待决策事项（含截止时间、Owner、潜在选项 A/B）。"),
    ("风险矩阵", "识别 5–8 条主要风险（技术 / 业务 / 组织 / 进度 / 数据），用 Markdown 表格列出：风险 | 影响 | 可能性 | 应对策略 | Owner。影响和可能性使用高/中/低。"),
    ("依赖与里程碑", "梳理关键交付物的依赖关系和里程碑节点，标出阻塞项。"),
    ("行业最佳实践对照", "结合知识库中同行业 / 同模块的实施经验，给出本项目可借鉴的 2–3 条做法和 1–2 条要规避的反例。引用具体出处文件名。"),
    ("下一步建议", "5–8 条可执行建议，每条注明 Owner、deadline、预期产出。区分 Quick Win（2 周内）与 Strategic（1–3 月）。"),
]


async def generate_insight(bundle_id: str, project_id: str):
    try:
        await _mark_bundle(bundle_id, "generating")
        ctx = await _gather_inputs(bundle_id, project_id, "insight")
        proj = ctx["project"]
        scope_line = f"项目名称：{proj.name}，客户：{proj.customer or '—'}" if proj else f"行业：{ctx['industry'] or '—'}"

        sections = []
        for title, question in INSIGHT_SECTIONS:
            prompt = f"""{scope_line}
行业：{ctx['industry'] or '—'}

【访谈记录】
{ctx['transcript']}

{f"【行业/客户研究 brief】{chr(10)}{ctx['industry_brief']}" if ctx.get('industry_brief') else ""}

{f"【知识库佐证】{chr(10)}{ctx['refs_text']}" if ctx['refs_text'] else ""}

{f"【联网检索】{chr(10)}{ctx['web_text']}" if ctx.get('web_text') else ""}

{f"【方法论】{chr(10)}{ctx['agent_prompt']}" if ctx['agent_prompt'] else ""}

{f"【启用技能】{chr(10)}{ctx['skill_text']}" if ctx['skill_text'] else ""}

{f"【已确认的项目 Brief（最权威素材，请优先采信，不要绕过）】{chr(10)}{ctx['project_brief']}" if ctx.get('project_brief') else ""}

本节主题：**{title}**
本节任务：{question}

【输出要求】
- 第一行就是 Bottom line 结论（加粗一句话）
- 必须用一张 Markdown 表格作为主要载体（除非"执行摘要"这种总览）
- 每个数据点末尾用 [访谈]/[知识库]/[行业 brief]/[Web]/[推断] 标注来源
- 字数 600–1200，禁止黑话
- 信息缺口写"信息缺失，建议在 Phase 1 第一周补访"，不要编造"""
            answer = await _llm_call(prompt, system=INSIGHT_SYSTEM, model=ctx["agent_model"], max_tokens=4000)
            answer = _strip_section_preamble(answer, title)
            sections.append(f"## {title}\n\n{answer}")

        report_date = date.today().strftime("%Y年%m月%d日")
        header_name = proj.name if proj else (ctx["industry"] or "—")
        md = (
            f"# {header_name} · 项目洞察报告\n\n"
            f"**生成日期**：{report_date}  \n"
            f"**客户**：{(proj.customer if proj else '—') or '—'}  \n"
            f"**行业**：{ctx['industry'] or '—'}\n\n---\n\n"
            + "\n\n---\n\n".join(sections)
        )

        await _mark_bundle(bundle_id, "done", content_md=md)
        await _mark_conversation(bundle_id, "done")
        logger.info("insight_generated", bundle_id=bundle_id, project_id=project_id)
    except Exception as e:
        logger.error("insight_failed", bundle_id=bundle_id, error=str(e)[:200])
        await _mark_bundle(bundle_id, "failed", error=str(e)[:500])
        await _mark_conversation(bundle_id, "failed")


# ── Kickoff PPTX（HTML 交付物，Claude 风格） ─────────────────────────────────

HTML_PPTX_SYSTEM = """你是一位 MBB 风格（McKinsey/BCG/Bain）咨询顾问，负责为 CRM 实施项目的启动会生成【可直接交付给甲方 C-level】的幻灯片。

【内容风格 — 严格遵守】
- 不要"好看"，要"专业"。深色文字、克制留白、表格 / 矩阵 / 图示优先
- 每页主标都是结论句（"基于现状诊断，主数据治理是 Phase 1 优先级最高的工作流"），副标可以是描述句
- 每页正文必须包含【至少 1 张表格 / 矩阵 / 图示】，不能是纯文字 bullet
- 每个论断必须带数字（行业 brief 中的"业界基准"也算）
- 信息源标签：每个数据点旁边用小字标 [访谈] / [知识库] / [行业 brief]；模型推断标 [推断]
- 不允许出现：赋能 / 抓手 / 闭环 / 链路 / 生态 / 一站式 / 全方位 / 数字化转型 / 全链路 / emoji
输出形态：一个完整自包含的 HTML 文件字符串，绝对不要 ```html 代码块围栏，直接从 <!DOCTYPE html> 开始。

【硬性规范】
- 16:9，每页固定 1280×720 像素，使用 <section class="slide"> 容器。
- 所有样式写在顶部 <style>，不使用 emoji，不引用外部图片。
- 字体栈："PingFang SC","Microsoft YaHei",-apple-system,"Helvetica Neue",sans-serif
- 色板（严格只用这几个）：主橙 #D96400，亮橙 #FB923C，墨黑 #1F2937，次灰 #4B5563，弱灰 #9CA3AF，分隔线 #E5E7EB，背景米白 #FAFAFA，纯白 #FFFFFF
- 12 栅格，边距 64px，默认行高 1.5。
- 多页之间用 CSS page-break 分页，便于浏览器打印 PDF：每个 .slide 设置 page-break-after: always。
- body 背景 #FAFAFA，slide 背景 #FFFFFF，阴影 0 4px 24px rgba(0,0,0,.06)。
- slide 设 overflow:hidden，宽 1280px 高 720px，居中 margin:24px auto。
- **不要写自己的翻页按钮 / 上下页快捷键**——平台会注入 .deck-nav 控制器。但你**应该**写 JS 来做图表渲染、SVG 交互、动画、数据可视化。

【视觉与 JS 增强（强烈鼓励，不是可选）】
- **inline `<svg>`** 是首选：甘特条/RACI 矩阵/陈vron 阶段流/网络拓扑/流程图都用 SVG path + text，比 CSS div 漂亮 10 倍
- **CDN 图表库可用**：Chart.js (`https://cdn.jsdelivr.net/npm/chart.js`)、d3 (`https://d3js.org/d3.v7.min.js`)、ECharts (`https://cdn.jsdelivr.net/npm/echarts`)——选一个用于柱状/饼/环形/雷达图
- **微交互**：CSS keyframes 入场动画（fadeIn/slideUp）、SVG path 动画 stroke-dasharray、悬停高亮
- **示例片段**：
  ```html
  <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
  <canvas id="riskChart" width="400" height="240"></canvas>
  <script>new Chart(document.getElementById('riskChart'),{type:'radar',data:{...}});</script>
  ```
- **SVG 甘特示例**：用 `<rect>` + `<text>`，里程碑用 `<polygon points="0,8 8,0 16,8 8,16">`（菱形）
- 一份"专业"PPT 应该至少包含 **3 种不同的 SVG/canvas 图**（甘特、雷达、矩阵），不能全是 div+border 做的伪表格

【字号层级】封面主标 52px / 副标 22px；页面标题 32px；小节标题 20px；正文 16px；图注/页脚 12px。

【必须包含的页（按顺序，缺失信息用"[待确认]"而不是编造）】
1. 封面：客户名 + "启动会" + 日期；左侧品牌色竖条；右下角"Fenxiao CRM · LTC 实施方法论"脚标
2. 议程：6 条编号目录，右侧橙色数字
3. 现状与挑战：左 2×2 矩阵（业务痛点×系统约束）+ 右摘要文案
4. 项目目标：3 条 SMART 目标卡片，每条含指标
5. 范围边界：In-scope / Out-of-scope 双列对照表
6. 方法论：LTC 阶段 chevron 流（机会→合同→交付→回款），当前阶段深橙，其他浅灰
7. 实施路径：甘特条（至少 4 阶段，双周粒度），行高 28px，橙条 + 里程碑菱形
8. 团队与治理：RACI 表（甲方 / 乙方 / 联合），单元格用 ● / ○ / 空
9. 风险与应对：表格（风险 | 影响 | 可能性 | 应对策略），影响/可能性用 高/中/低 彩色标签
10. 资源与投入：人天分布柱状条 + 关键角色卡
11. 下一步 Next Step：本周 / 下周 Action Items，每条含 Owner + deadline + 橙色圆点

【视觉元件实现提示】
- chevron 用 clip-path: polygon(0 0, calc(100% - 20px) 0, 100% 50%, calc(100% - 20px) 100%, 0 100%, 20px 50%);
- 2×2 矩阵：优先 inline SVG（坐标轴 + 象限标签 + 数据点），次选 grid-template-columns: 1fr 1fr
- 甘特条：优先 inline SVG（rect + 菱形里程碑 polygon + 周次 text 轴），不要用 div+absolute 凑
- RACI 表：可用 table，但配色方块 (●○) 用 SVG circle 比 unicode 字符更精致
- 风险矩阵 / 资源分布：建议 Chart.js radar / bar，CDN 加载

【文案规范】
- 每页标题必须是结论句（"基于现状诊断，优先打通三大主数据"），不是描述句（"现状分析"）
- 不要咨询黑话（赋能/抓手/闭环/生态/链路）
- 每个主张至少 1 个数字（若访谈没有数据，标 [待确认]）

【输出】直接输出完整 HTML 字符串，不要任何解释、不要 markdown 围栏、不要前后语。从 <!DOCTYPE html> 起到 </html> 止。"""


PPTGEN_PYTHON_SYSTEM = """你是一位 MBB 风格（McKinsey/BCG/Bain）咨询顾问 + Python 工程师，要为 CRM 实施项目的启动会生成【可直接交付给甲方 C-level】的真·PPTX 文件。

**输出形态：纯 Python 代码字符串**——可被 `python xxx.py` 直接执行，最终在当前目录生成 `out.pptx`。绝对不要 ```python 围栏，绝对不要任何解释文字、markdown，只输出代码。

【内容风格 — 严格遵守】
- 不要"好看"，要"专业"。深色文字、克制留白、表格 / 矩阵 / 图示优先
- 每页主标都是结论句（"基于现状诊断，主数据治理是 Phase 1 优先级最高的工作流"），副标可以是描述句
- 每页正文必须包含【至少 1 张表格 / 矩阵 / 图示】，不能是纯文字 bullet
- 每个论断必须带数字（行业 brief 中的"业界基准"也算）
- 信息源标签：每个数据点旁边用小字标 [访谈] / [知识库] / [行业 brief]；模型推断标 [推断]
- 不允许出现：赋能 / 抓手 / 闭环 / 链路 / 生态 / 一站式 / 全方位 / 数字化转型 / 全链路 / emoji
- 文案简练：每条 bullet ≤ 25 汉字；标题 ≤ 22 汉字

【硬性栅格 — 严格遵守，否则会重叠！】
- 幻灯片 16:9：宽 13.333"，高 7.5"
- **左右半区分隔（关键）**：
  - 左半区任何元素的 x + w ≤ **6.4**（即 x 起点 0.6，宽度上限 5.8）
  - 右半区任何元素的 x ≥ **6.9**（即 x 起点 6.9，宽度上限 5.8，到 12.7）
  - "左 X + 右 Y" 布局时，左侧元素绝对不能跨过 6.4，右侧不能跨过 6.9 往左
- **垂直分区**：标题 y ∈ [0.4, 1.3]，正文 y ∈ [1.5, 6.7]，页脚 y ∈ [6.9, 7.3]
- 每个元素之间至少留 0.2" gap，绝对不允许 bbox 重叠
- 在写代码前先在注释里声明每页的 x/y 分区表，再绘制

【硬性技术规范 — Python 代码层面】
- 必须包含：`import matplotlib; matplotlib.use("Agg"); import matplotlib.pyplot as plt`
- 必须设 matplotlib 中文字体（容器里已装 Noto CJK，但 .ttc 在 matplotlib font_manager 里只注册了 JP face，JP 字形包含全部汉字，对中文显示完全 OK）：
  ```
  plt.rcParams['font.sans-serif'] = ['Noto Sans CJK JP', 'Noto Sans CJK SC', 'WenQuanYi Zen Hei', 'DejaVu Sans']
  plt.rcParams['axes.unicode_minus'] = False
  ```
  **必须把 'Noto Sans CJK JP' 放第一位**，否则 matplotlib 找不到字体降级到 DejaVu，中文渲染成方框。
- 必须使用 `from pptx import Presentation; from pptx.util import Inches, Pt, Emu; from pptx.dml.color import RGBColor; from pptx.enum.shapes import MSO_SHAPE; from pptx.enum.text import PP_ALIGN, MSO_ANCHOR`
- 16:9 设置：`prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)`
- 全部使用 `prs.slide_layouts[6]`（blank layout）
- 色板（用 RGBColor，严格只用这几个）：
  - BRAND_ORANGE = RGBColor(0xD9, 0x64, 0x00)
  - LIGHT_ORANGE = RGBColor(0xFB, 0x92, 0x3C)
  - INK = RGBColor(0x1F, 0x29, 0x37)
  - SUB = RGBColor(0x4B, 0x55, 0x63)
  - MUTED = RGBColor(0x9C, 0xA3, 0xAF)
  - DIVIDER = RGBColor(0xE5, 0xE7, 0xEB)
  - BG = RGBColor(0xFA, 0xFA, 0xFA)
  - WHITE = RGBColor(0xFF, 0xFF, 0xFF)
- 字号：封面主标 44pt / 副标 20pt；页面标题 28pt；小节标题 18pt；正文 14pt；图注 10pt
- **中文字体统一用 ZH_FONT = "Microsoft YaHei"（黑体系，跨 Win/Mac 都有）**——所有 r.font.name 必须显式赋 ZH_FONT，禁止用 PingFang SC（Windows 上没有会回落到衬线字体）。matplotlib 图内中文字体无关 ZH_FONT，由 plt.rcParams['font.sans-serif'] 决定（已设 Noto Sans CJK JP）。

【色调统一 — 强制】**全篇只能出现上面 8 个色板色，不能出现任何其他颜色**（包括 python-pptx / matplotlib 的默认蓝、绿、红、紫）：
- **表格（add_table）禁止使用默认蓝色样式**。创建表后必须遍历每个 cell 显式设置：
  - 表头行：`cell.fill.solid(); cell.fill.fore_color.rgb = BRAND_ORANGE`；文字 WHITE 加粗
  - 数据行（奇/偶）：背景 WHITE 或 BG，文字 INK / SUB
  - 边框：python-pptx 表格边框需要操作 XML，可改用 add_rect 拼网格替代 add_table（更可控）
- **matplotlib 颜色严格映射**到品牌色（先把 RGBColor 转 hex 字符串：`BRAND_HEX = "#D96400"` 等，matplotlib 用这些 hex）：
  - 柱状图主色：`color='#D96400'`（BRAND_ORANGE）；分组配色用 `['#D96400', '#FB923C', '#9CA3AF']`
  - 雷达图：填充 `'#D96400'` 透明度 0.25，描边 `'#D96400'` 实线宽 2
  - 网格线：`'#E5E7EB'`（DIVIDER）；轴标签 `'#4B5563'`（SUB）；标题 `'#1F2937'`（INK）
  - 禁止使用 matplotlib 默认 'tab10' / 'C0..C9' 配色循环（默认蓝 #1f77b4 与品牌色冲突）
- **chevron / shape 高亮**：当前阶段 BRAND_ORANGE，未激活阶段 DIVIDER（不要用浅灰、深灰等其他灰）
- **里程碑 DIAMOND**：填充 BRAND_ORANGE，描边 WHITE
- **任何 matplotlib spines / ticks 颜色**：`ax.spines[s].set_color('#9CA3AF')`，避免黑色硬边

【matplotlib 图表 — 关键防扭曲规则】
- **figsize 与 add_picture 必须等比例**！算法：picture 占位框 width_in × height_in（PPT 中 inch），figsize=(width_in, height_in)（matplotlib 中 inch），dpi=150
- **保留比例 add_picture**：`slide.shapes.add_picture(buf, Inches(x), Inches(y), width=Inches(W))`，**只传 width，不传 height**——让 PPT 按图自动算高度。或同时传 width 和 height，但保证比例和 figsize 完全一致
- **极坐标雷达**：`fig, ax = plt.subplots(figsize=(4.5, 4.5), subplot_kw={'projection':'polar'})`，必须正方形！add_picture 也用正方形 width=height
- 背景透明：`fig.patch.set_alpha(0); ax.patch.set_alpha(0)`
- 保存：`fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', transparent=True)`
- savefig 后必须 `plt.close(fig)` 释放
- 中文允许出现在 matplotlib 内（已装 CJK 字体），但**所有数据标签优先放外部 add_textbox**，图表内只放最少必要文本

【11 页结构 — 严格按此顺序，缺数据用"[待确认]"，不要编造】
1. 封面：左侧品牌色竖条（宽 0.5"）；中央客户名 + "项目启动会"主标 + 日期；右下角"Fenxiao CRM · LTC 实施方法论"脚标
2. 议程：6 条编号目录
3. 现状与挑战：**左 2×2 矩阵（x ∈ [0.6, 6.4]）+ 右摘要文字（x ∈ [6.9, 12.7]）**。矩阵示例数据结构：
   ```
   # 2×2 数据：每个 quadrant 是 {title, bullets[]}
   quadrants = [
       {"title": "业务痛点·内部", "bullets": ["销售流程在线率<30% [访谈]", "客户信息分散 [访谈]"]},
       {"title": "业务痛点·外部", "bullets": ["商机转化周期 45 天 [行业 brief]", "客户流失率 12% [推断]"]},
       {"title": "系统约束·当前", "bullets": ["6 套独立系统 [访谈]", "数据孤岛 [访谈]"]},
       {"title": "系统约束·未来", "bullets": ["无统一权限 [推断]", "无 BI 集成 [访谈]"]},
   ]
   # 4 个 cell 排成 2×2，cell 宽 (5.8-0.2)/2=2.8，cell 高 (5.0-0.2)/2=2.4
   for idx, q in enumerate(quadrants):
       row, col = idx // 2, idx % 2
       cx = 0.6 + col * 3.0
       cy = 1.6 + row * 2.6
       add_rect(slide, cx, cy, 2.8, 2.4, fill=BG, line=DIVIDER)
       add_text(slide, cx+0.15, cy+0.1, 2.5, 0.4, q["title"], size=14, bold=True, color=BRAND_ORANGE)
       for k, b in enumerate(q["bullets"]):
           add_text(slide, cx+0.15, cy+0.6+k*0.5, 2.5, 0.5, "• "+b, size=11, color=SUB)
   ```
4. 项目目标：3 张 SMART 目标卡，每张宽 ≈ 4.0，间距 0.2，y ∈ [1.6, 6.4]
5. 范围边界：左 In-scope 表 + 右 Out-of-scope 表（**严格守左右半区**）
6. 方法论：LTC 4 阶段 chevron 流（机会 → 合同 → 交付 → 回款），**当前阶段必须是第一阶段"机会"**（i==0 高亮 BRAND_ORANGE，其他 DIVIDER），因为是 kickoff
7. 实施路径：甘特图，rect 横条 + DIAMOND 里程碑，至少 4 阶段
8. 团队与治理：RACI 表（首列任务名 + 4 列角色 × 至少 8 行）
9. 风险与应对：**左 雷达图（width=height=4.5"，正方形！）+ 右 风险表**
10. 资源与投入：**左 柱状图（width=6.0", height=4.0", figsize=(6.0, 4.0)）+ 右 角色卡**
11. 下一步：本周 / 下周双列 Action Items

【代码模板骨架（必须保留这些工具函数和初始化，11 页正文请你自己展开）】
```
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
plt.rcParams['font.sans-serif'] = ['Noto Sans CJK JP', 'Noto Sans CJK SC', 'WenQuanYi Zen Hei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BRAND_ORANGE = RGBColor(0xD9, 0x64, 0x00)
LIGHT_ORANGE = RGBColor(0xFB, 0x92, 0x3C)
INK = RGBColor(0x1F, 0x29, 0x37)
SUB = RGBColor(0x4B, 0x55, 0x63)
MUTED = RGBColor(0x9C, 0xA3, 0xAF)
DIVIDER = RGBColor(0xE5, 0xE7, 0xEB)
BG = RGBColor(0xFA, 0xFA, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
# matplotlib hex 等价
HEX_BRAND = "#D96400"; HEX_LIGHT = "#FB923C"; HEX_INK = "#1F2937"
HEX_SUB = "#4B5563"; HEX_MUTED = "#9CA3AF"; HEX_DIVIDER = "#E5E7EB"
ZH_FONT = "Microsoft YaHei"  # 黑体系字体，Win/Mac 都默认安装；避免 PingFang 在 Windows 上失效

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, font=ZH_FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return tb

def add_rect(slide, x, y, w, h, *, fill=BRAND_ORANGE, line=None, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s

# ……此处展开全部 11 页 ……

prs.save("out.pptx")
```

【自检清单 — 写完每一页后默念，然后写下一页】
1. 这一页有没有元素的 bbox 越过 x=6.4 和 x=6.9 的中线（如果是左右布局）？
2. 这一页有没有 matplotlib 图，如果有，figsize 和 add_picture 的 inch 比例是否完全一致？
3. 这一页的标题是不是结论句 + 数字？
4. 表格 / 图示是否真出现（不能纯文字 bullet）？
5. 这一页所有颜色是否都来自 8 色色板（无默认蓝、绿、红、紫等乱入）？matplotlib 是否显式传 color/edgecolor，没有依赖默认配色？表格 cell 是否显式 fill.solid + RGBColor，没有用默认蓝表样式？

【最终输出】只输出可直接执行的 Python 源码，从 import 行起，到 `prs.save("out.pptx")` 止。不要 ``` 围栏，不要任何 markdown 或解释文字。"""


async def generate_kickoff_pptx(bundle_id: str, project_id: str):
    try:
        await _mark_bundle(bundle_id, "generating")
        ctx = await _gather_inputs(bundle_id, project_id, "kickoff_pptx")
        proj = ctx["project"]

        if proj:
            customer = proj.customer or ""
            kickoff_date_str = proj.kickoff_date.strftime("%Y年%m月%d日") if proj.kickoff_date else "待定"
            scope_block = f"""项目名称：{proj.name}
客户：{customer or "未填写"}
行业：{ctx['industry'] or "未填写"}
启动日期：{kickoff_date_str}
实施模块：{", ".join(proj.modules or []) or "未填写"}
项目描述：{proj.description or "无"}"""
            title_name = proj.name
        else:
            customer = ""
            kickoff_date_str = "待定"
            scope_block = f"行业：{ctx['industry'] or '—'}\n（没有具体项目信息，按行业共性展开）"
            title_name = ctx["industry"] or "行业启动会"

        prompt = f"""{scope_block}

【访谈记录】
{ctx['transcript']}

{f"【行业/客户研究 brief（务必把这部分变成 PPT 内容，不要只读一遍）】{chr(10)}{ctx['industry_brief']}" if ctx.get('industry_brief') else ""}

{f"【知识库佐证】{chr(10)}{ctx['refs_text']}" if ctx['refs_text'] else ""}

{f"【联网检索】{chr(10)}{ctx['web_text']}" if ctx.get('web_text') else ""}

{f"【方法论/风格要求】{chr(10)}{ctx['agent_prompt']}" if ctx['agent_prompt'] else ""}

{f"【启用技能（PPT 骨架 / 版式 / 文案规范）】{chr(10)}{ctx['skill_text']}" if ctx['skill_text'] else ""}

{f"【已确认的项目 Brief（最权威素材，PPT 各页内容必须基于此展开，不要绕过）】{chr(10)}{ctx['project_brief']}" if ctx.get('project_brief') else ""}

请输出可执行的 Python 源代码，使用 python-pptx 生成 11 页启动会 .pptx，最终 prs.save("out.pptx")。每页都要有表格/矩阵/图示，不要纯文字 bullet。"""

        # PPT 生成默认走小米 mimo-v2-pro；max_tokens=None 不设上限（mimo 先 think 再出代码，需要给足空间）
        from services.pptx_codeexec import execute_pptx_code, strip_python_fences

        pptx_model = ctx["agent_model"] or "mimo-v2-pro"
        code_raw = await _llm_call(prompt, system=PPTGEN_PYTHON_SYSTEM, model=pptx_model, max_tokens=None, timeout=1200.0)
        code = strip_python_fences(code_raw)
        if "prs.save" not in code or "from pptx" not in code:
            raise RuntimeError(f"模型未输出有效 python-pptx 代码：开头 200 字符 = {code[:200]!r}")

        logger.info("pptx_code_received", bundle_id=bundle_id, code_len=len(code))
        pptx_bytes = await execute_pptx_code(code, timeout=180.0)

        pptx_key = f"outputs/{bundle_id}/kickoff.pptx"
        _minio_put(pptx_key, pptx_bytes,
                   "application/vnd.openxmlformats-officedocument.presentationml.presentation")

        # 同时保留生成代码作为调试快照（次要文件，下载链接仍指向 .pptx）
        code_key = f"outputs/{bundle_id}/kickoff_gen.py"
        _minio_put(code_key, code.encode("utf-8"), "text/x-python; charset=utf-8")

        md = f"# {customer or title_name} · 启动会 PPT\n\n" \
             f"**生成日期**：{date.today().strftime('%Y-%m-%d')}  \n" \
             f"**客户**：{customer or '—'}  \n" \
             f"**行业**：{ctx['industry'] or '—'}\n\n" \
             f"> 真 .pptx 文件已生成（{len(pptx_bytes)//1024} KB），点击下载后用 PowerPoint / Keynote / WPS 直接打开编辑。"

        await _mark_bundle(bundle_id, "done", content_md=md, file_key=pptx_key)
        await _mark_conversation(bundle_id, "done")
        logger.info("pptx_generated", bundle_id=bundle_id, project_id=project_id, format="pptx",
                    size=len(pptx_bytes))
    except Exception as e:
        # httpx.ReadTimeout 等异常的 str(e) 可能为空，统一拼上类型名便于排查
        msg = f"{type(e).__name__}: {str(e) or '(no message)'}"
        logger.error("pptx_failed", bundle_id=bundle_id, error=msg[:500])
        await _mark_bundle(bundle_id, "failed", error=msg[:500])
        await _mark_conversation(bundle_id, "failed")


# ── Kickoff HTML（HTML 版幻灯片，浏览器内播放） ─────────────────────────────────

async def generate_kickoff_html(bundle_id: str, project_id: str):
    """旧版 HTML 启动会 PPT：模型直出 HTML，浏览器即开即看。与 generate_kickoff_pptx 并存。"""
    try:
        await _mark_bundle(bundle_id, "generating")
        ctx = await _gather_inputs(bundle_id, project_id, "kickoff_html")
        proj = ctx["project"]

        if proj:
            customer = proj.customer or ""
            kickoff_date_str = proj.kickoff_date.strftime("%Y年%m月%d日") if proj.kickoff_date else "待定"
            scope_block = f"""项目名称：{proj.name}
客户：{customer or "未填写"}
行业：{ctx['industry'] or "未填写"}
启动日期：{kickoff_date_str}
实施模块：{", ".join(proj.modules or []) or "未填写"}
项目描述：{proj.description or "无"}"""
            title_name = proj.name
        else:
            customer = ""
            scope_block = f"行业：{ctx['industry'] or '—'}\n（没有具体项目信息，按行业共性展开）"
            title_name = ctx["industry"] or "行业启动会"

        prompt = f"""{scope_block}

【访谈记录】
{ctx['transcript']}

{f"【行业/客户研究 brief（务必把这部分变成 PPT 内容，不要只读一遍）】{chr(10)}{ctx['industry_brief']}" if ctx.get('industry_brief') else ""}

{f"【知识库佐证】{chr(10)}{ctx['refs_text']}" if ctx['refs_text'] else ""}

{f"【联网检索】{chr(10)}{ctx['web_text']}" if ctx.get('web_text') else ""}

{f"【方法论/风格要求】{chr(10)}{ctx['agent_prompt']}" if ctx['agent_prompt'] else ""}

{f"【启用技能（PPT 骨架 / 版式 / 文案规范）】{chr(10)}{ctx['skill_text']}" if ctx['skill_text'] else ""}

{f"【已确认的项目 Brief（最权威素材，PPT 各页内容必须基于此展开，不要绕过）】{chr(10)}{ctx['project_brief']}" if ctx.get('project_brief') else ""}

请输出完整自包含 HTML，11 页 16:9 幻灯片（1280×720），从 <!DOCTYPE html> 起到 </html> 止。"""

        html_model = ctx["agent_model"]  # 走 doc_generation 路由
        html = await _llm_call(prompt, system=HTML_PPTX_SYSTEM, model=html_model, max_tokens=None, timeout=900.0)
        # 去 ``` 围栏
        s = (html or "").strip()
        if s.startswith("```"):
            nl = s.find("\n")
            s = s[nl + 1:] if nl >= 0 else s
            if s.endswith("```"):
                s = s[:-3]
            s = s.strip()
        if "<!DOCTYPE" not in s and "<html" not in s.lower():
            raise RuntimeError(f"模型未输出有效 HTML：开头 200 字符 = {s[:200]!r}")

        html_key = f"outputs/{bundle_id}/kickoff.html"
        _minio_put(html_key, s.encode("utf-8"), "text/html; charset=utf-8")

        md = f"# {customer or title_name} · 启动会 PPT（HTML 版）\n\n" \
             f"**生成日期**：{date.today().strftime('%Y-%m-%d')}  \n" \
             f"**客户**：{customer or '—'}  \n" \
             f"**行业**：{ctx['industry'] or '—'}\n\n" \
             f"> HTML 幻灯片已生成（{len(s)//1024} KB），点击预览即可全屏播放，浏览器打印即可导出 PDF。"

        await _mark_bundle(bundle_id, "done", content_md=md, file_key=html_key)
        await _mark_conversation(bundle_id, "done")
        logger.info("html_generated", bundle_id=bundle_id, project_id=project_id, format="html",
                    size=len(s))
    except Exception as e:
        msg = f"{type(e).__name__}: {str(e) or '(no message)'}"
        logger.error("html_failed", bundle_id=bundle_id, error=msg[:500])
        await _mark_bundle(bundle_id, "failed", error=msg[:500])
        await _mark_conversation(bundle_id, "failed")


def _strip_section_preamble(answer: str, section_title: str) -> str:
    """去掉模型在 section answer 里重复输出的整篇标题、封面元信息表、与本节标题重名的 H1/H2。"""
    import re as _re
    s = (answer or "").strip()
    if s.startswith("```"):
        first_nl = s.find("\n")
        if first_nl >= 0:
            s = s[first_nl + 1:]
        if s.endswith("```"):
            s = s[:-3]
        s = s.strip()
    lines = s.split("\n")
    meta_keys = ("客户名称", "项目代号", "报告日期", "编写人", "版本号", "生成日期")
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if not line or line in ("---", "***"):
            i += 1
            continue
        if _re.match(r"^#{1,2}\s+", line):
            i += 1
            continue
        # 封面元信息表：连续 | 开头的行，包含"客户名称/项目代号/报告日期"等
        if line.startswith("|"):
            j = i
            while j < len(lines) and lines[j].strip().startswith("|"):
                j += 1
            block = "\n".join(lines[i:j])
            if any(k in block for k in meta_keys):
                i = j
                continue
            break
        break
    return "\n".join(lines[i:]).strip()


def _strip_html_fences(raw: str) -> str:
    """剥掉模型可能输出的 ```html / ``` 围栏。"""
    s = raw.strip()
    if s.startswith("```"):
        # 去掉首行围栏
        first_nl = s.find("\n")
        if first_nl >= 0:
            s = s[first_nl + 1:]
        if s.endswith("```"):
            s = s[:-3]
    return s.strip()


SLIDESHOW_CSS = """
/* 幻灯片切换模式 */
@media screen {
  body { background:#0f1115 !important; min-height:100vh; display:flex; flex-direction:column; align-items:center; padding:32px 0 96px !important; }
  .slide { display:none !important; margin:0 auto !important; }
  .slide.is-active { display:block !important; }
  .deck-nav {
    position:fixed; left:0; right:0; bottom:0; height:64px;
    background:rgba(20,22,28,.92); backdrop-filter:blur(8px);
    display:flex; align-items:center; justify-content:center; gap:14px;
    color:#fff; font-family:"PingFang SC","Microsoft YaHei",sans-serif; font-size:13px;
    border-top:1px solid #1f2937; z-index:9999;
  }
  .deck-nav button {
    background:#D96400; color:#fff; border:none; border-radius:8px; padding:8px 16px;
    font-size:13px; cursor:pointer; font-weight:600;
  }
  .deck-nav button:disabled { background:#4B5563; cursor:not-allowed; }
  .deck-nav .deck-counter { min-width:80px; text-align:center; color:#9CA3AF; font-variant-numeric:tabular-nums; }
  .deck-nav .deck-print { background:transparent; border:1px solid #4B5563; color:#9CA3AF; }
  .deck-nav .deck-hint { color:#6B7280; font-size:12px; margin-left:8px; }
}
@media print {
  body { background:#fff !important; padding:0 !important; }
  .slide { display:block !important; page-break-after:always; box-shadow:none !important; }
  .deck-nav { display:none !important; }
}
"""

SLIDESHOW_JS = """
(function(){
  var slides = Array.prototype.slice.call(document.querySelectorAll('.slide'));
  if (!slides.length) return;
  var i = 0;
  function go(n){ i = Math.max(0, Math.min(slides.length - 1, n)); slides.forEach(function(s, idx){ s.classList.toggle('is-active', idx === i); }); update(); slides[i].scrollIntoView({behavior:'instant',block:'center'}); }
  function update(){
    var c = document.getElementById('deck-counter'); if (c) c.textContent = (i+1) + ' / ' + slides.length;
    var p = document.getElementById('deck-prev'); if (p) p.disabled = (i === 0);
    var n = document.getElementById('deck-next'); if (n) n.disabled = (i === slides.length - 1);
  }
  document.addEventListener('keydown', function(e){
    if (e.key === 'ArrowRight' || e.key === 'PageDown' || e.key === ' ') { e.preventDefault(); go(i+1); }
    else if (e.key === 'ArrowLeft' || e.key === 'PageUp') { e.preventDefault(); go(i-1); }
    else if (e.key === 'Home') { go(0); } else if (e.key === 'End') { go(slides.length-1); }
  });
  window.__deckGo = go;
  window.__deckStep = function(d){ go(i + d); };
  window.__deckPrint = function(){ window.print(); };
  go(0);
})();
"""

SLIDESHOW_NAV_HTML = """
<nav class="deck-nav">
  <button id="deck-prev" type="button" onclick="window.__deckStep(-1)">← 上一页</button>
  <span class="deck-counter" id="deck-counter">1 / 1</span>
  <button id="deck-next" type="button" onclick="window.__deckStep(1)">下一页 →</button>
  <button class="deck-print" type="button" onclick="window.__deckPrint()">打印 / 导出 PDF</button>
  <span class="deck-hint">键盘 ← → 切换</span>
</nav>
"""


def _inject_slideshow(html: str) -> str:
    """在生成的 HTML 中注入幻灯片切换样式与脚本。
    - 在 </head> 前追加 SLIDESHOW_CSS；如果没有 </head>，包一层。
    - 在 </body> 前追加导航条 + 脚本。
    """
    import re
    s = html

    # 1) 注入 CSS：优先放进现有 <style>；其次 </head> 前；都没有则不动（fallback HTML 自带样式）
    style_tag = f"<style id=\"deck-overrides\">{SLIDESHOW_CSS}</style>"
    if "</head>" in s.lower():
        # 用大小写不敏感方式插入
        s = re.sub(r"</head>", style_tag + "</head>", s, count=1, flags=re.IGNORECASE)
    elif "<body" in s.lower():
        s = re.sub(r"<body([^>]*)>", style_tag + r"<body\1>", s, count=1, flags=re.IGNORECASE)
    else:
        s = style_tag + s

    nav_and_script = SLIDESHOW_NAV_HTML + f"<script>{SLIDESHOW_JS}</script>"
    if "</body>" in s.lower():
        s = re.sub(r"</body>", nav_and_script + "</body>", s, count=1, flags=re.IGNORECASE)
    else:
        s = s + nav_and_script
    return s


def _fallback_html(title: str, customer: str, kickoff_date: str, raw: str) -> str:
    """模型没吐正确 HTML 时的兜底：把原文塞进一个简单壳里，至少可下载。"""
    import html as _h
    body = _h.escape(raw)
    return f"""<!DOCTYPE html>
<html lang="zh-CN"><head><meta charset="UTF-8"><title>{_h.escape(customer or title)} 启动会</title>
<style>body{{font-family:"PingFang SC","Microsoft YaHei",sans-serif;background:#FAFAFA;color:#1F2937;padding:48px;line-height:1.7}}
h1{{color:#D96400}} pre{{white-space:pre-wrap;background:#fff;padding:24px;border-radius:8px;border:1px solid #E5E7EB}}</style>
</head><body><h1>{_h.escape(customer or title)} 启动会</h1><p>日期：{_h.escape(kickoff_date)}</p>
<p style="color:#9CA3AF;font-size:12px">（HTML 解析失败，以下为模型原始输出）</p>
<pre>{body}</pre></body></html>"""


# ── Helpers ───────────────────────────────────────────────────────────────────



def _build_docx(title: str, md: str) -> bytes:
    from docx import Document as DocxDoc

    doc = DocxDoc()
    doc.add_heading(title, level=0)
    for line in md.splitlines():
        line = line.rstrip()
        if line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
        elif line:
            doc.add_paragraph(line)
    buf = io.BytesIO(); doc.save(buf)
    return buf.getvalue()


def _minio_put(key: str, data: bytes, content_type: str):
    from minio import Minio
    mc = Minio(settings.minio_endpoint, access_key=settings.minio_user, secret_key=settings.minio_password, secure=False)
    mc.put_object(settings.minio_bucket, key, io.BytesIO(data), len(data), content_type=content_type)
